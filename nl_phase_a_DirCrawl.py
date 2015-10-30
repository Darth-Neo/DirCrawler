#
# Crawl a directory for documents and pull out the text
#

from nl_lib.Logger import *
logger = setupLogging(__name__)
logger.setLevel(INFO)

from nl_lib.Constants import *
from nl_lib.Concepts import Concepts

import nltk

# by importing the textract module, you will then enable its use
if False:
    from textract import *
    TEXTRACT = True
    logger.info(u"Using textract parser")

else: # ImportError, msg:
    TEXTRACT = False
    logger.info(u"Using custom parser : ") # %s" % msg)

import openxmllib
from pptx import Presentation
from n_docx import opendocx, getdocumenttext
import xlrd
from pyPdf import PdfFileReader

from traceback import format_exc
from unidecode import unidecode
from BeautifulSoup import BeautifulSoup

stop.append(u"information")
stop.append(u"member")

import unicodedata

class DirCrawl(object):
    documentsConceptsFile = u"documents.p"
    documentsConcepts = None

    wordsConceptsFile = u"words.p"
    wordsConcepts = None
    
    def __init__(self):
        self.documentsConcepts = Concepts(u"DocumentConcepts", u"Documents")
        self.wordsConcepts = Concepts(u"WordConcepts", u"Words")

    def getDocumentsConcepts(self):
        return self.documentsConcepts

    def getWordsConcepts(self):
        return self.wordsConcepts

    def _saveConcepts(self):
        Concepts.saveConcepts(self.documentsConcepts, self.documentsConceptsFile)
        Concepts.saveConcepts(self.wordsConcepts, self.wordsConceptsFile)

    def _getOpenXmlText(self, filename, c):
        logger.debug(u"OpenXmlText: %s" % filename)

        try:
            doc = openxmllib.openXmlDocument(path=filename)
            logger.debug(u"%s\n" % doc.allProperties)

            ap = c.addConceptKeyType(u"allProperties", u"PROPERTIES")
            for x in doc.allProperties:
                logger.debug(u"cp %s:%s" % (x, doc.allProperties[x]))
                ap.addConceptKeyType(doc.allProperties[x], x)
        except Exception, msg:
            logger.error(u"%s" % msg)

    def _getPDFText(self, filename, d):
        logger.debug(u"filename: %s" % filename)
        newparatextlist = list()

        try:
            pdfDoc = PdfFileReader(file(filename, u"rb"))

            pdfDict = pdfDoc.getDocumentInfo()

            for x in pdfDict.keys():
                d.addConceptKeyType(x[1:], pdfDict[x])

            # c.logConcepts()

            for page in pdfDoc.pages:
                text = page.extractText()
                if not isinstance(text, str):
                    unicodedata.normalize(u'NFKD', text).encode(u'ascii', u'ignore')

                logger.debug(u"PDF : %s" % text)

                newparatextlist.append(text + u". ")

            return newparatextlist

        except Exception, msg:
            logger.error(u"%s" % msg)

    def _getPPTXText(self, filename):
        logger.debug(u"filename: %s" % filename)

        newparatextlist = list()

        try:
            prs = Presentation(filename)

            for slide in prs.slides:
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            logger.debug(u"PPTX : %s" % run.text)
                            if run.text is not None:
                                if not isinstance(run.text, str):
                                    run.text = unicode(run.text)

                                newparatextlist.append(run.text + u". ")
        except Exception, msg:
            logger.error(u"%s" % msg)

        return newparatextlist

    def _getXLSText(self, filename):
        logger.debug(u"filename: %s" % filename)

        newparatextlist = list()

        try:
            workbook = xlrd.open_workbook(filename)

            # sheet = "Specific Requirements"
            # worksheet = workbook.sheet_by_name(sheet)

            CellTypes = [u"Empty", u"Text", u"Number", u"Date", u"Boolean", u"Error", u"Blank"]

            for worksheet_name in workbook.sheet_names():
                worksheet = workbook.sheet_by_name(worksheet_name)
                num_rows = worksheet.nrows - 1
                num_cells = worksheet.ncols - 1
                curr_row = -1

                while curr_row < num_rows:
                    curr_row += 1
                    row = worksheet.row(curr_row)
                    logger.debug(u'Row: %s' % curr_row)
                    curr_cell = -1
                    while curr_cell < num_cells:
                        curr_cell += 1
                        # Cell Types: 0=Empty, 1=Text, 2=Number, 3=Date, 4=Boolean, 5=Error, 6=Blank
                        cell_type = worksheet.cell_type(curr_row, curr_cell)
                        cell_value = worksheet.cell_value(curr_row, curr_cell)
                        if cell_type == 1:
                            if not isinstance(cell_value, str):
                                cell_value = unicode(cell_value)
                            logger.debug(u"XLXS : %s" % cell_value)
                            newparatextlist.append(cell_value + u". ")
        except Exception, msg:
            logger.error(u"%s" % msg)

        return newparatextlist

    def _getDOCXText(self, filename):
        logger.debug(u"filename: %s" % filename)

        document = opendocx(filename)
        # Fetch all the text out of the document we just created

        paratextlist = getdocumenttext(document)

        # Make explicit unicode version
        newparatextlist = []
        for paratext in paratextlist:
            if not isinstance(paratext, str):
                paratext = unicode(paratext)
            logger.debug(u"DOCX : %s" % paratext)
            newparatextlist.append(paratext + u". ")
            
        return newparatextlist

    def _getTXT(self, filename):
        logger.debug(u"filename: %s" % filename)

        listText = list()

        try:
            soup = BeautifulSoup(open(filename))
        except Exception, msg:
            logger.warn(u"%s" % msg)
            return None

        logger.debug(u"Soup: %s" % soup)

        st = soup.text

        for line in st.split(os.linesep):
            logger.debug(u"TXT : %s" % line)
            listText.append(line)

        return listText

    def _getConcepts(self, fname, d, w, DOCX=True, XLSX=True, PDF=True, PPTX=True, MISC=True):
        listText = list()
        text = None

        if TEXTRACT is True:
            try:
                txt = process(str(fname)).encode(u"utf8", errors=u"replace")
                text = txt.decode(u"ascii", errors=u"replace")
                text = unidecode(txt)
                if text is not None or len(text) != 0:
                    listText = text
            except Exception, msg:
                logger.error(u"%s" % msg)
            logger.info(u"++parsing : %s[%d]" % (fname, len(listText)))
        else:
            txt = str(fname).encode(u"utf8", errors=u"replace")
            text = txt.decode(u"ascii", errors=u"replace")

            if fname[-5:] == u".docx" and DOCX is True:
                listText = self._getDOCXText(fname)
                self._getOpenXmlText(fname, d)
                logger.info(u"+c %s" % fname)
            elif fname[-5:] == u".pptx" and PPTX is True:
                listText = self._getPPTXText(fname)
                self._getOpenXmlText(fname, d)
                logger.info(u"++Parsing = %s" % fname)
            elif fname[-5:] == u" .xlsx" and XLSX is True:
                listText = self._getXLSText(fname)
                self._getOpenXmlText(fname, d)
                logger.info(u"++Parsing = %s" % fname)
            elif fname[-4:] == u".pdf" and PDF is True:
                listText = self._getPDFText(fname, d)
                logger.info(u"++Parsing = %s" % fname)
            elif MISC is True:
                if fname[-4:] in (u".txt", u".xml", u".htm", u".csv"):
                    listText = self._getTXT(fname)
                    logger.info(u"++Parsing = %s" % fname)
                elif fname[-5:] in (u".html", u".WSDL"):
                    listText = self._getTXT(fname)
                    logger.info(u"++Parsing = %s" % fname)
                elif fname[-4:] in (u".xml"):
                    listText = self._getTXT(fname)
                    logger.info(u"++Parsing = %s" % fname)
            else:
                logger.warn(u"Unsupported file type: %s" % (fname))

        if listText is not None:
            if isinstance(listText, list):
                for t in listText:
                    if t is not None:
                        try:
                            # c = unicode(t).strip()
                            # sentence = t.encode('ascii', errors="ignore").strip()
                            logger.debug(u"%s:Text : %s" % (type(t), t))
                            w = d.addConceptKeyType(t, u"Text")
                            self._addWords(w, t)
                        except Exception, msg:
                            logger.error(u"%s" % msg)
            else:
                w = d.addConceptKeyType(listText, u"Text")
                self._addWords(w, listText)

        return listText

    def checkWordCaps(self, w):
        s = str()

    def _addWords(self, words, sentence):
            cleanSentence = u' '.join([word for word in sentence.split() if word not in stop])

            logger.debug(u"cs:%s" % cleanSentence)

            for word, pos in nltk.pos_tag(nltk.wordpunct_tokenize(cleanSentence)):
                logger.debug(u"Word: " + word + u" POS: " + pos)
                c = words.addConceptKeyType(word, u"WORD")
                c.addConceptKeyType(pos, u"POS")

    def _checkFile(self, fname):
        logger.debug(u"filename: %s" % fname)

        fna = unidecode(fname)
        # fna = fn.encode("ascii", errors="ignore")

        d = Concepts(fna, u"Document")
        w = Concepts(fna, u"Document")

        listText = self._getConcepts(fna, d, w)

        if listText is not None and len(listText) == 0:
            logger.warn(u"%s has no text" % fna)
            return 0
        else:
            self.documentsConcepts.addConcept(d)
            self.wordsConcepts.addConcept(w)

            return 1
                
    def searchSubDir(self, subdir):
        numFilesParsed = 0
        for root, dirs, files in os.walk(subdir, topdown=True):
            for name in files:
                nameFile = os.path.join(root, name)
                numFilesParsed += self._checkFile(nameFile)

        self._saveConcepts()

        return numFilesParsed


def test_dirCrawl():

    # Set the directory you want to start from
    rootDir = os.getcwd() + os.sep + u"testdata"

    numFilesParsed = 0

    dc = DirCrawl()

    numFilesParsed = dc.searchSubDir(rootDir)

    logger.info(u"Documents Parsed = %d" % numFilesParsed)

    return dc

if __name__ == u'__main__':

    test_dirCrawl()
