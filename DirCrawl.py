# Import the os module, for the os.walk function
import sys
import os
import csv

from nl_lib import Logger
logger = Logger.setupLogging(__name__)

from nl_lib import Constants
from nl_lib import Concepts 
#from nl_lib import PeopleText
#from nl_lib import Tokens
#from nltk.corpus import stopwords
#from nltk import tokenize

from pptx import Presentation


def getPPTX(filename):
    path_to_presentation = "C:\\Users\\morrj140\\Dev\\GitRepository\\DirCrawler\\Introduction to GBTS - v1.1.pptx"
    prs = Presentation(filename)

    newparatextlist = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_textframe:
                continue
            for paragraph in shape.textframe.paragraphs:
                for run in paragraph.runs:
                    newparatextlist.append(run.text.encode("utf-8"))
    return newparatextlist

from docx import opendocx, getdocumenttext

documentsConcepts = Concepts.Concepts("DocumentConcepts", "Documents")

def getText(filename):
    document = opendocx(filename)
    
    # Fetch all the text out of the document we just created
    paratextlist = getdocumenttext(document)

    # Make explicit unicode version
    newparatextlist = []
    for paratext in paratextlist:
        newparatextlist.append(paratext.encode("utf-8"))
        
    return newparatextlist

def getConcepts(fname, d):
    logger.debug("filename: %s" % fname)

    if fname[-5:] == ".docx" or fname[-4:] == ".doc":
        listText = getText(fname)
    elif fname[-5:] == ".pptx":
        listText = getPPTX(fname)
                
    for t in listText:
        d.addConceptKeyType(t.strip(), "Text")

    return listText

def checkFile(fname):
    logger.debug("filename: %s" % fname)
            
    try:
        d = documentsConcepts.addConceptKeyType(fname, "Document")
        getConcepts(fname, d)
    except:
        logger.warn("File could not be parsed : %s" % fname)
            
def searchSubDir(subdir):
    i = 0
    for root, dirs, files in os.walk(rootDir, topdown=False):
        for name in files:
            nameFile = os.path.join(root, name)
            i += 1
            logger.info(nameFile)
            checkFile(nameFile)

        for name in dirs:
            nameFile = os.path.join(root, name)
            i += 1
            logger.info(nameFile)
            checkFile(nameFile)

    logger.info("=====Complete=====")
    logger.info("Checked %d documents" % i)
  
if __name__ == '__main__':  
    # Set the directory you want to start from
    #rootDir = '.'
    #rootDir = 'C:\\Users\\morrj140\\Documents\\'
    rootDir = "C:\\Users\\morrj140\\Documents\\System Architecture\\Accovia"
    
    searchSubDir(rootDir)

    Concepts.Concepts.saveConcepts(documentsConcepts, "documents.p")
            
